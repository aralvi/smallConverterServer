from fastapi import UploadFile, File, HTTPException, status
from typing import Optional
from pydantic import BaseModel
from app.schemas.conversion_schema import ConversionMeta
import uuid
import os
import io
from app.utils.get_file_mime import file_type
from app.utils.file_to_base64 import file_to_base64
from app.utils.clear_temp import clear_temp
from decouple import config
# for image to image & image to pdf
from PIL import Image
# for pdf to image
import fitz
# for pdf to docx
from pdf2docx import Converter as pdf_2_docx
import aiofiles
# for pdf to excel
import camelot as camelot
# for pdf to pptx
from tqdm import trange
from pptx import Presentation, action
from pptx.util import Cm
# for word, excel, powerpoint to pdf
from win32com import client as office_client


if (config("PLATFORM", cast=str) == "windows"):
    current_path = config("CURRENT_PATH", cast=str)+"temp\\"
else:
    current_path = config("CURRENT_PATH", cast=str)+"temp/"


class Conversions(BaseModel):
    file: str
    filename: str
    filetype: str


async def convert_file(conversion_type: str, file: UploadFile, meta: ConversionMeta) -> Optional[Conversions]:
    # Clearing Temporary Files
    clear_temp(current_path)
    file_name = os.path.splitext(file.filename)[
        0]+"-converted"+f".{meta['tgt']}"
    file_type = os.path.splitext(file.filename)[1]
    # ^^ Gives us a tuple with filename at index 0 and type at 1
    # Now converting it into bytes
    file_bytes = file.file.read()
    file.file.close()

    match conversion_type:
        case "imgtoimg":
            result = img_to_img(file_bytes, meta["tgt"])
        case "imgtopdf":
            result = img_to_pdf(file_bytes)
        case "pdftoimg":
            result = pdf_to_img(file_bytes, meta["tgt"])
        case "pdftodocx":
            result = await pdf_to_docx(file_bytes, meta["tgt"])
        case "pdftoxlsx":
            result = await pdf_to_xlsx(file_bytes, meta["tgt"])
        case "pdftopptx":
            result = await pdf_to_pptx(file_bytes)
        case "docxtopdf":
            result = await docx_to_pdf(file_bytes, file_type)
        case "xlsxtopdf":
            result = await xlsx_to_pdf(file_bytes, file_type)
        case "pptxtopdf":
            result = await pptx_to_pdf(file_bytes, file_type)
        case _:
            return None

    if not result:
        return None
    # elif not result[0]:
    #     raise HTTPException(
    #         status_code=status.HTTP_400_BAD_REQUEST,
    #         detail={
    #             "message": result[1]
    #         }f
    #     )

    return {"file": result["file"], "filename": file_name, "filetype": result["filetype"]}


def img_to_img(file_bytes: bytes = File(), tgt: str = "png") -> dict:
    if tgt == "jpg":
        tgt = "jpeg"
    # loading image
    image = Image.open(io.BytesIO(file_bytes))
    output_buffer = io.BytesIO()
    # Converting
    if image.mode != "RGB":
        image.convert("RGB").save(output_buffer, format=tgt.upper())
    else:
        image.save(output_buffer, format=tgt.upper())
    image.close()
    # finalizing
    # type = file_type(f"temp/{name}.{tgt}")
    b64 = file_to_base64(output_buffer.getvalue())
    return {"file": b64, "filetype": "type"}


def img_to_pdf(file_bytes: bytes = File()) -> dict:
    # loading image
    image = Image.open(io.BytesIO(file_bytes))
    output_buffer = io.BytesIO()
    # Converting
    if image.mode != "RGB":
        image.convert("RGB")
    image.save(output_buffer, "PDF", resolution=100.0)
    image.close()
    # finalizing
    # type = file_type(f"temp/{name}.pdf")
    b64 = file_to_base64(output_buffer.getvalue())
    return {"file": b64, "filetype": "type"}


def pdf_to_img(file_bytes: bytes = File(), tgt: str = "png") -> dict:
    if tgt == "jpg":
        tgt = "jpeg"
    # loading document
    doc = fitz.Document(stream=io.BytesIO(file_bytes))
    # converting
    page = doc.load_page(0)
    pix = page.get_pixmap()
    output_buffer = pix.pil_tobytes(format=tgt.upper(), optimize=True)
    doc.close()
    # finalizing
    # type = file_type(f"temp/{name}.{tgt}")
    b64 = file_to_base64(output_buffer)
    return {"file": b64, "filetype": "type"}


async def pdf_to_docx(file_bytes: bytes = File(), tgt: str = "docx") -> dict:
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}.pdf", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading document
    pdf = pdf_2_docx(f"temp/{name}.pdf")
    # converting
    pdf.convert(f"temp/{name}.{tgt}")
    pdf.close()
    # deleting saved doc
    if os.path.exists(f"temp/{name}.pdf"):
        os.remove(f"temp/{name}.pdf")
    # finalizing
    type = file_type(f"temp/{name}.{tgt}")
    b64 = file_to_base64(f"temp/{name}.{tgt}")
    return {"file": b64, "filetype": type}


async def pdf_to_xlsx(file_bytes: bytes = File(), tgt: str = "xlsx") -> dict:
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}.pdf", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading document
    # reading table from page 2
    table = camelot.read_pdf(f"temp/{name}.pdf", pages="1", flavor="stream")
    # converting
    table[0].df.to_excel(f"temp/{name}.{tgt}")
    # deleting saved doc
    if os.path.exists(f"temp/{name}.pdf"):
        os.remove(f"temp/{name}.pdf")
    # finalizing
    type = file_type(f"temp/{name}.{tgt}")
    b64 = file_to_base64(f"temp/{name}.{tgt}")
    return {"file": b64, "filetype": "type"}


async def pdf_to_pptx(file_bytes: bytes = File(), tgt: str = "pptx") -> dict:
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}.pdf", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading document
    doc = fitz.open(f"temp/{name}.pdf")
    page_count = doc.page_count
    # transformation matrix: slide to pixmap
    zoom = 300 / 72
    matrix = fitz.Matrix(zoom, zoom, 0)
    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    # configure presentation aspect ratio
    page = doc.load_page(0)
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)
    # create page iterator
    page_iter = range(0, page_count)
    # iterate over slides
    for page_no in page_iter:
        page = doc.load_page(page_no)

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)
        image_data = pixmap.tobytes(output='PNG')
        image_file = io.BytesIO(image_data)

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(
            image_file, left, top, height=prs.slide_height)
    # saving presentation
    prs.save(f"temp/{name}.{tgt}")
    # closing document
    doc.close()
    # deleting saved doc
    if os.path.exists(f"temp/{name}.pdf"):
        os.remove(f"temp/{name}.pdf")
    # finalizing
    type = file_type(f"temp/{name}.{tgt}")
    b64 = file_to_base64(f"temp/{name}.{tgt}")
    return {"file": b64, "filetype": "type"}


async def docx_to_pdf(file_bytes: bytes = File(), src: str = "") -> dict | tuple:
    if src == "":
        return (None, "Specify the source file type, doc or docx in this case")
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}{src}", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading word
    word = office_client.Dispatch('Word.Application')
    # reading document
    doc = word.Documents.Open(f"{current_path}{name}{src}")
    # converting
    doc.SaveAs(f"{current_path}{name}.pdf", FileFormat=17)
    # closing document and word
    doc.Close()
    word.Quit()
    # deleting saved doc
    if os.path.exists(f"temp/{name}{src}"):
        os.remove(f"temp/{name}{src}")
    # finalizing
    type = file_type(f"temp/{name}.pdf")
    b64 = file_to_base64(f"temp/{name}.pdf")
    return {"file": b64, "filetype": type}


async def xlsx_to_pdf(file_bytes: bytes = File(), src: str = "") -> dict | tuple:
    if src == "":
        return (None, "Specify the source file type, xls or xlsx in this case")
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}{src}", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading excel
    excel = office_client.Dispatch("Excel.Application")
    # reading xlsx
    sheets = excel.Workbooks.Open(
        f"{current_path}{name}{src}")
    work_sheets = sheets.Worksheets[0]
    # Convert to pdf
    work_sheets.ExportAsFixedFormat(
        0, f"{current_path}{name}.pdf")
    sheets.Close()
    # Quitting excel
    excel.Quit()
    # deleting saved excel file
    if os.path.exists(f"temp/{name}{src}"):
        os.remove(f"temp/{name}{src}")
    # finalizing
    type = file_type(f"temp/{name}.pdf")
    b64 = file_to_base64(f"temp/{name}.pdf")
    return {"file": b64, "filetype": type}


async def pptx_to_pdf(file_bytes: bytes = File(), src: str = "") -> dict | tuple:
    if src == "":
        return (None, "Specify the source file type, ppt or pptx in this case")
    name = str(uuid.uuid4().hex)
    # saving the file
    async with aiofiles.open(f"temp/{name}{src}", "wb") as saved_file:
        await saved_file.write(file_bytes)
    # loading powerpoint
    powerpoint = office_client.Dispatch("Powerpoint.Application")
    # reading pptx
    deck = powerpoint.Presentations.Open(
        f"{current_path}{name}{src}", WithWindow=False)
    # Convert to pdf
    # 32 denoting ppt_to_pdf
    deck.SaveAs(
        f"{current_path}{name}.pdf", 32)
    # complete list here https://learn.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype
    # stackoverflow answer https://stackoverflow.com/questions/31487478/how-to-convert-a-pptx-to-pdf-using-python
    deck.Close()
    # Quitting powerpoint
    powerpoint.Quit()
    # deleting saved powerpoint file
    if os.path.exists(f"temp/{name}{src}"):
        os.remove(f"temp/{name}{src}")
    # finalizing
    type = file_type(f"temp/{name}.pdf")
    b64 = file_to_base64(f"temp/{name}.pdf")
    return {"file": b64, "filetype": type}
